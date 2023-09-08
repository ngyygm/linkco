import os
import json
import docx
import pptx
import PyPDF2
import codecs
import chardet
import pandas as pd


def read_csv_file(file_path: str) -> dict:
    """
    读取CSV文件，并将数据保存到字典中

    Args:
        file_path (str): CSV文件路径

    Returns:
        dict: 保存了CSV文件数据的字典
    """
    try:
        data = pd.read_csv(file_path, encoding='gbk')
    except:
        data = pd.read_csv(file_path, encoding='utf-8')

    return data.to_dict(orient='list')


def read_xlsx_file(file_path: str, sheet_name: str = None) -> dict:
    """
    读取xlsx文件，并将数据保存到字典中

    Args:
        file_path (str): xlsx文件路径
        sheet_name (str): Sheet名称，可选

    Returns:
        dict: 保存了xlsx文件数据的字典
    """
    data = pd.read_excel(file_path, sheet_name=sheet_name)
    if isinstance(data, dict):
        return data
    else:
        return data.to_dict(orient='list')


def save_dict_to_csv_file(data: dict, save_path: str) -> None:
    """
    将字典数据保存为CSV文件

    Args:
        data (dict): 需要保存的字典数据
        save_path (str): 保存路径

    Returns:
        None
    """
    df = pd.DataFrame(data)
    df.to_csv(save_path, index=False, encoding='utf-8-sig')


def save_dict_to_xlsx_file(data: dict, save_path: str, sheet_name: str = 'Sheet1') -> None:
    """
    将字典数据保存为xlsx文件

    Args:
        data (dict): 需要保存的字典数据
        save_path (str): 保存路径
        sheet_name (str): Sheet名称，可选

    Returns:
        None
    """
    df = pd.DataFrame(data)
    df.to_excel(save_path, index=False, encoding='utf-8-sig', sheet_name=sheet_name)


def save_dict_to_json_file(data: dict, save_path: str) -> None:
    """
    将字典数据保存为JSON文件

    Args:
        data (dict): 需要保存的字典数据
        file_path (str): 保存路径

    Returns:
        None
    """
    with open(save_path, "w") as file:
        json.dump(data, file, indent=4)


def read_pdf_file(file_path: str) -> str:
    """
    Read a PDF file and extract the text content.

    Args:
        file_path (str): Path to the PDF file.

    Returns:
        str: Text content extracted from the PDF.
    """
    with open(file_path, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        text = ''
        for page in reader.pages:
            text += page.extract_text()
    return text


def read_txt_file(file_path: str) -> str:
    """
    Read a TXT file and extract the text content.

    Args:
        file_path (str): Path to the TXT file.

    Returns:
        str: Text content extracted from the TXT file.
    """

    with open(file_path, 'rb') as file:
        raw_data = file.read()
        result = chardet.detect(raw_data)
        encoding = result['encoding']

    with codecs.open(file_path, 'r', encoding=encoding) as file:
        text = file.read()

    return text


def read_docx_file(file_path: str) -> str:
    """
    Read a DOCX file and extract the text content.

    Args:
        file_path (str): Path to the DOCX file.

    Returns:
        str: Text content extracted from the DOCX file.
    """
    doc = docx.Document(file_path)
    paragraphs = []
    for para in doc.paragraphs:
        paragraphs.append(para.text)
    return '\n'.join(paragraphs)


def read_pptx_file(file_path: str) -> str:
    """
    Read a PPTX file and extract the text content.

    Args:
        file_path (str): Path to the PPTX file.

    Returns:
        str: Text content extracted from the PPTX file.
    """
    prs = pptx.Presentation(file_path)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + ' '
            text += '\n'
    return text


def read_file(file_path: str) -> str:
    """
    Read a file and extract its text content.

    Args:
        file_path (str): Path to the file.

    Returns:
        str: Text content extracted from the file.
    """
    file_extension = os.path.splitext(file_path)[1].lower()

    if file_extension == ".pdf":
        return read_pdf_file(file_path)
    elif file_extension == ".txt":
        return read_txt_file(file_path)
    elif file_extension == ".docx" or file_extension == ".doc":
        return read_docx_file(file_path)
    elif file_extension == ".csv":
        return read_csv_file(file_path)
    elif file_extension == ".xlsx" or file_extension == ".xls":
        return read_xlsx_file(file_path)
    elif file_extension == ".pptx" or file_extension == ".ppt":
        return read_pptx_file(file_path)
    else:
        try:
            return read_txt_file(file_path)
        except Exception as e:
            print(e)
            return "Unsupported file format."